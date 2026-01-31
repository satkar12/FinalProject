import os
import pytesseract
import nltk
import torch
import re
from fastapi import HTTPException
from typing import Dict
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
from pptx import Presentation
from transformers import T5Tokenizer, T5ForConditionalGeneration
from nltk.tokenize import sent_tokenize

# ---------------- INIT ----------------
nltk.download("punkt")

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

app = FastAPI(title="QuickPrep AI Summarizer")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------- LOAD MODEL ----------------
print("Loading T5 model...")
# ---------------- LOAD MODEL ----------------
print("Loading T5 model from Hugging Face...")

# 1. Replace this with your actual Hugging Face Repo ID (e.g., "username/quickprep-t5")
model_repo_id = "satkar12/quickprep-t5" 

# 2. Transformers will automatically download and cache the model for you
tokenizer = T5Tokenizer.from_pretrained(model_repo_id)
model = T5ForConditionalGeneration.from_pretrained(model_repo_id)

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
model.to(device)

print(f"Model loaded on {device}.")

# ---------------- FILE READERS ----------------
def read_pdf(path):
    text = ""
    reader = PdfReader(path)
    for page in reader.pages:
        if page.extract_text():
            text += page.extract_text() + "\n"
    return text

def read_scanned_pdf(path):
    text = ""
    images = convert_from_path(path)
    for img in images:
        text += pytesseract.image_to_string(img) + "\n"
    return text

# ---------------- TEXT NORMALIZATION ----------------
def normalize_text(text):
    text = re.sub(r"\n+", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()

def prepare_text_for_summary(text, min_words=60):
    words = text.split()
    if len(words) < min_words:
        pass
    return text

# ---------------- HEADING DETECTION (PDF) ----------------
def is_heading_pdf(line: str) -> bool:
    return (
        line.isupper()
        or re.match(r"^\d+(\.\d+)*\s+.+", line)
        or (len(line.split()) <= 6 and not line.endswith("."))
    )

def split_pdf_by_headings(text: str) -> Dict[str, str]:
    lines = text.split("\n")
    sections = {}
    current_heading = "Introduction"
    sections[current_heading] = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        if is_heading_pdf(line):
            current_heading = line
            sections[current_heading] = []
        else:
            sections[current_heading].append(line)

    return {k: " ".join(v) for k, v in sections.items()}

# ---------------- HEADING DETECTION (PPTX) ----------------
def extract_pptx_blocks(path):
    prs = Presentation(path)
    blocks = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                blocks.append(p.text.strip())

    return blocks

def split_pptx_by_headings(path):
    blocks = extract_pptx_blocks(path)
    sections = {}
    current_heading = "Introduction"
    sections[current_heading] = []

    for text in blocks:
        if not text:
            continue
        if len(text.split()) <= 6 and not text.endswith("."):
            current_heading = text
            sections[current_heading] = []
        else:
            sections[current_heading].append(text)

    return {k: " ".join(v) for k, v in sections.items()}

# ---------------- TEXT CHUNKING ----------------
def chunk_text(text, max_len=800):
    sentences = sent_tokenize(text)
    chunks, current = [], ""

    for s in sentences:
        if len(current) + len(s) <= max_len:
            current += s + " "
        else:
            chunks.append(current.strip())
            current = s + " "

    if current:
        chunks.append(current.strip())

    return chunks

# ---------------- SUMMARIZATION ----------------
def summarize_text(text: str) -> str:
    text = normalize_text(text)
    text = prepare_text_for_summary(text)

    chunks = chunk_text(text)
    summaries = []

    # ----- First pass -----
    for chunk in chunks:
        prompt = (
            "summarize:" + chunk
        )

        inputs = tokenizer.encode(
            prompt,
            return_tensors="pt",
            truncation=True,
            max_length=512
        ).to(device)

        output = model.generate(
            inputs,
            max_length=150,
            min_length=30,
            num_beams=4,
            length_penalty=1.0,
            no_repeat_ngram_size=3,
            repetition_penalty=2.5,
            early_stopping=True,
        )

        summaries.append(
            tokenizer.decode(output[0], skip_special_tokens=True)
        )

    

    final_summary = tokenizer.decode(output[0], skip_special_tokens=True)

    final_summary = re.sub(r"[\n•\-]+", " ", final_summary)
    final_summary = re.sub(r"\s+", " ", final_summary).strip()

    return final_summary

# ---------------- API ENDPOINTS ----------------
@app.post("/summarize-pdf")
async def summarize_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")

    path = os.path.join(UPLOAD_DIR, file.filename)
    with open(path, "wb") as f:
        f.write(await file.read())

    text = read_pdf(path)
    if not text.strip():
        text = read_scanned_pdf(path)

    if not text.strip():
        raise HTTPException(status_code=400, detail="Could not extract text.")

    sections = split_pdf_by_headings(text)
    results = []

    for heading, body in sections.items():
        if body.strip():
            results.append({
                "heading": heading,
                "summary": summarize_text(body)
            })

    return {"sections": results}

@app.post("/summarize-pptx")
async def summarize_pptx(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only PPTX files are supported.")

    path = os.path.join(UPLOAD_DIR, file.filename)
    with open(path, "wb") as f:
        f.write(await file.read())

    sections = split_pptx_by_headings(path)
    results = []

    for heading, body in sections.items():
        if body.strip():
            results.append({
                "heading": heading,
                "summary": summarize_text(body)
            })

    return {"sections": results}
